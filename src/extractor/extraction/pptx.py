"""PowerPoint (.pptx) extraction."""

from __future__ import annotations

from pathlib import Path

from extractor.extraction.base import ExtractionError


def extract(path: Path) -> str:
    from pptx import Presentation

    try:
        presentation = Presentation(str(path))
    except Exception as e:
        raise ExtractionError(f"python-pptx could not open file: {e}") from e

    slides_text: list[str] = []
    for i, slide in enumerate(presentation.slides, start=1):
        parts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    run_text = "".join(run.text for run in paragraph.runs)
                    if run_text.strip():
                        parts.append(run_text)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells]
                    if any(cells):
                        parts.append(" | ".join(cells))
        if parts:
            slides_text.append(f"--- Slide {i} ---\n" + "\n".join(parts))

    text = "\n\n".join(slides_text)
    if not text.strip():
        raise ExtractionError("No extractable text found in .pptx")
    return text
