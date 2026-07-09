"""Shared test fixtures.

Sample documents are generated on the fly (via the same libraries the
extractors use to read them) rather than committed as static binary files —
keeps the repo free of binary fixture blobs and keeps fixture content
colocated with the tests that use it.
"""

from __future__ import annotations

from pathlib import Path

import pytest


@pytest.fixture
def make_txt(tmp_path: Path):
    def _make(content: str = "Hello from a plain text file.\nSecond line.") -> Path:
        p = tmp_path / "sample.txt"
        p.write_text(content, encoding="utf-8")
        return p

    return _make


@pytest.fixture
def make_docx(tmp_path: Path):
    def _make(paragraphs: list[str] | None = None) -> Path:
        import docx

        document = docx.Document()
        for para in paragraphs or ["Hello from a Word document.", "Second paragraph."]:
            document.add_paragraph(para)
        p = tmp_path / "sample.docx"
        document.save(str(p))
        return p

    return _make


@pytest.fixture
def make_pptx(tmp_path: Path):
    def _make(slide_texts: list[str] | None = None) -> Path:
        from pptx import Presentation

        presentation = Presentation()
        layout = presentation.slide_layouts[1]
        for text in slide_texts or ["Hello from a slide.", "Second slide content."]:
            slide = presentation.slides.add_slide(layout)
            slide.shapes.title.text = text
        p = tmp_path / "sample.pptx"
        presentation.save(str(p))
        return p

    return _make


@pytest.fixture
def make_xlsx(tmp_path: Path):
    def _make(rows: list[list[str]] | None = None) -> Path:
        import openpyxl

        workbook = openpyxl.Workbook()
        sheet = workbook.active
        for row in rows or [["name", "value"], ["alpha", "1"], ["beta", "2"]]:
            sheet.append(row)
        p = tmp_path / "sample.xlsx"
        workbook.save(str(p))
        return p

    return _make


@pytest.fixture
def make_csv(tmp_path: Path):
    def _make(rows: list[list[str]] | None = None) -> Path:
        import csv

        p = tmp_path / "sample.csv"
        with p.open("w", newline="", encoding="utf-8") as f:
            writer = csv.writer(f)
            for row in rows or [["name", "value"], ["alpha", "1"], ["beta", "2"]]:
                writer.writerow(row)
        return p

    return _make


@pytest.fixture
def make_html(tmp_path: Path):
    def _make(body: str = "<h1>Title</h1><p>Hello from HTML.</p>") -> Path:
        p = tmp_path / "sample.html"
        p.write_text(
            f"<html><head><style>body{{color:red}}</style></head><body>{body}</body></html>"
        )
        return p

    return _make


@pytest.fixture
def make_eml(tmp_path: Path):
    def _make(subject: str = "Test Subject", body: str = "Hello from an email body.") -> Path:
        from email.message import EmailMessage

        msg = EmailMessage()
        msg["From"] = "sender@example.com"
        msg["To"] = "recipient@example.com"
        msg["Subject"] = subject
        msg.set_content(body)
        p = tmp_path / "sample.eml"
        p.write_bytes(bytes(msg))
        return p

    return _make


@pytest.fixture
def make_pdf(tmp_path: Path):
    def _make(text: str = "Hello from a PDF document.") -> Path:
        from fpdf import FPDF

        pdf = FPDF()
        pdf.add_page()
        pdf.set_font("Helvetica", size=12)
        pdf.cell(text=text)
        p = tmp_path / "sample.pdf"
        pdf.output(str(p))
        return p

    return _make
